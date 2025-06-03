import streamlit as st
import pandas as pd
from typing import List, Dict, Tuple, Optional
import tempfile
import os
from pathlib import Path
import base64
from io import BytesIO
import json
import time
import textwrap

# Document processing imports
try:
    import fitz  # PyMuPDF for PDF
    import docx  # python-docx for Word documents
    from pptx import Presentation  # python-pptx for PowerPoint
    import mammoth  # For better DOC/DOCX text extraction
    from PIL import Image
    import pytesseract  # For OCR if needed
except ImportError as e:
    st.error(f"Missing required library: {e}")
    st.stop()

# Video generation imports
try:
    # MoviePy 2.2.1 compatible imports
    from moviepy import VideoFileClip, AudioFileClip, ImageClip, TextClip, ColorClip, CompositeVideoClip, concatenate_videoclips
    import cv2
    import numpy as np
    from gtts import gTTS  # Google Text-to-Speech
    
    # Test MoviePy installation
    test_clip = ColorClip(size=(100, 100), color=(0, 0, 0), duration=0.1)
    test_clip.close()
    
except ImportError as e:
    st.error(f"Missing video processing library: {e}")
    st.markdown("""
    ### 🔧 Installation Instructions
    
    For MoviePy 2.2.1 compatibility:
    
    **Recommended installation:**
    ```bash
    pip install "moviepy==2.2.1"
    pip install imageio-ffmpeg
    ```
    
    **Alternative methods:**
    ```bash
    # Method 1 - With conda
    conda install -c conda-forge moviepy
    
    # Method 2 - Latest version
    pip install moviepy
    ```
    
    After installation, restart your Python environment/kernel.
    """)
    st.stop()
except Exception as e:
    st.error(f"MoviePy installation issue: {e}")
    st.markdown("""
    ### Additional Requirements
    
    MoviePy also requires:
    - **FFmpeg** (for video processing)
    - **ImageIO** (for image handling)
    
    **To install FFmpeg:**
    
    **Windows:**
    ```bash
    # Using chocolatey
    choco install ffmpeg
    
    # Or download from https://ffmpeg.org/download.html
    ```
    
    **macOS:**
    ```bash
    brew install ffmpeg
    ```
    
    **Linux (Ubuntu/Debian):**
    ```bash
    sudo apt update
    sudo apt install ffmpeg
    ```
    
    **Alternative: Install imageio-ffmpeg**
    ```bash
    pip install imageio-ffmpeg
    ```
    """)
    st.stop()

class DocumentProcessor:
    """Handle document text and image extraction"""
    
    @staticmethod
    def extract_from_pdf(file_path: str) -> Tuple[List[str], List[Image.Image]]:
        """Extract text and images from PDF"""
        doc = fitz.open(file_path)
        texts = []
        images = []
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            
            # Extract text
            text = page.get_text()
            if text.strip():
                texts.append(text.strip())
            
            # Extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    img_data = pix.tobytes("png")
                    img_pil = Image.open(BytesIO(img_data))
                    images.append(img_pil)
                pix = None
        
        doc.close()
        return texts, images
    
    @staticmethod
    def extract_from_docx(file_path: str) -> Tuple[List[str], List[Image.Image]]:
        """Extract text and images from DOCX"""
        doc = docx.Document(file_path)
        texts = []
        images = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                texts.append(paragraph.text.strip())
        
        # Extract images from document parts
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_data = rel.target_part.blob
                    img_pil = Image.open(BytesIO(img_data))
                    images.append(img_pil)
                except Exception:
                    continue
        
        return texts, images
    
    @staticmethod
    def extract_from_pptx(file_path: str) -> Tuple[List[str], List[Image.Image]]:
        """Extract text and images from PPTX"""
        prs = Presentation(file_path)
        texts = []
        images = []
        
        for slide in prs.slides:
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract images
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        img_data = image.blob
                        img_pil = Image.open(BytesIO(img_data))
                        images.append(img_pil)
                    except Exception:
                        continue
            
            if slide_text:
                texts.append(" ".join(slide_text))
        
        return texts, images

class TextExpander:
    """Expand and enhance text content"""
    
    @staticmethod
    def expand_text(text: str, target_length: int = 100) -> str:
        """Expand text to meet minimum length requirements"""
        words = text.split()
        if len(words) >= target_length:
            return text
        
        # Simple expansion by adding explanatory phrases
        expanded_parts = []
        for sentence in text.split('.'):
            if sentence.strip():
                expanded_parts.append(sentence.strip())
                # Add explanatory content based on context
                if any(keyword in sentence.lower() for keyword in ['process', 'method', 'approach']):
                    expanded_parts.append("This involves several key steps and considerations")
                elif any(keyword in sentence.lower() for keyword in ['result', 'outcome', 'conclusion']):
                    expanded_parts.append("The implications of this are significant for understanding the overall concept")
                elif any(keyword in sentence.lower() for keyword in ['important', 'key', 'main']):
                    expanded_parts.append("This point deserves special attention and careful analysis")
        
        return '. '.join(expanded_parts) + '.'

class ProgressTracker:
    """Track and display progress for video generation"""
    
    def __init__(self, total_steps: int):
        self.total_steps = max(1, total_steps)  # Ensure total_steps is at least 1
        self.current_step = 0
        self.progress_bar = st.progress(0)
        self.status_text = st.empty()
        self.detail_text = st.empty()
        
    def update(self, step_name: str, detail: str = ""):
        """Update progress bar and status"""
        self.current_step += 1
        # Ensure progress never exceeds 1.0
        progress = min(self.current_step / self.total_steps, 1.0)
        self.progress_bar.progress(progress)
        self.status_text.text(f"Step {min(self.current_step, self.total_steps)}/{self.total_steps}: {step_name}")
        if detail:
            self.detail_text.text(detail)
        
    def set_progress(self, progress_value: float, step_name: str, detail: str = ""):
        """Set progress directly with a value between 0.0 and 1.0"""
        # Ensure progress is within valid bounds
        progress = max(0.0, min(progress_value, 1.0))
        self.progress_bar.progress(progress)
        self.status_text.text(step_name)
        if detail:
            self.detail_text.text(detail)
        
    def complete(self):
        """Mark as complete"""
        self.progress_bar.progress(1.0)
        self.status_text.text("✅ Video generation complete!")
        self.detail_text.text("")

class VideoGenerator:
    """Generate video from matched content"""
    
    def __init__(self, output_dir: str):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.progress_tracker = None
    
    def set_progress_tracker(self, tracker: ProgressTracker):
        """Set progress tracker for this generator"""
        self.progress_tracker = tracker
    
    def create_text_clip(self, text: str, duration: float, size: Tuple[int, int] = (1920, 1080)) -> TextClip:
        """Create text clip with improved styling"""
        if self.progress_tracker:
            self.progress_tracker.update("Creating text overlay", f"Processing: {text[:50]}...")
        
        # Create base text clip with better visibility - MoviePy 2.2.1 compatible
        try:
            # Try common system fonts in order of preference
            fonts_to_try = ['Arial', 'DejaVu-Sans', 'Liberation-Sans']
            font_used = None
            wrapped_text = textwrap.fill(text, width=70)  # Wrap text to approx. 70 chars per line
            for font in fonts_to_try:
                try:
                    txt_clip = TextClip(
                        text=wrapped_text,
                        font_size=40,
                        color='white', 
                        font=font,
                        stroke_color='black',
                        stroke_width=2,
                        method='caption',
                        size=(1600, None)  # Width defined; height auto-adjusted
                    ).with_duration(duration)
                    font_used = font
                    break
                except:
                    continue
            
            # If no specific font works, use default
            if font_used is None:
                txt_clip = TextClip(
                    text=wrapped_text,
                    font_size=40,
                    color='white',
                    stroke_color='black',
                    stroke_width=2,
                    method='caption',
                    size=(1600, None)
                ).with_duration(duration)
                
        except Exception as e:
            # Fallback to simplest TextClip creation
            txt_clip = TextClip(
                text=wrapped_text,
                font_size=30,
                color='white'
            ).with_duration(duration)
        
        return txt_clip
    
    def create_image_clip_with_animation(self, image: Image.Image, duration: float, 
                                       animation_type: str = "zoom", size: Tuple[int, int] = (1920, 1080)) -> ImageClip:
        """Create image clip with animation effects using resize transforms"""
        if self.progress_tracker:
            self.progress_tracker.update("Adding image animation", f"Animation: {animation_type}")
        
        # Save image temporarily
        temp_path = self.output_dir / f"temp_img_{id(image)}.png"
        
        # Resize image to full screen while maintaining aspect ratio
        img_ratio = image.width / image.height
        screen_ratio = size[0] / size[1]
        
        if img_ratio > screen_ratio:
            # Image is wider - fit to height
            new_height = size[1]
            new_width = int(new_height * img_ratio)
        else:
            # Image is taller - fit to width
            new_width = size[0]
            new_height = int(new_width / img_ratio)
        
        # Resize and center crop to full screen
        resized_img = image.resize((new_width, new_height), Image.LANCZOS)
        
        # Create full screen image by center cropping
        left = (new_width - size[0]) // 2
        top = (new_height - size[1]) // 2
        right = left + size[0]
        bottom = top + size[1]
        
        if left < 0 or top < 0:
            # If image is smaller than screen, create background and paste
            full_img = Image.new('RGB', size, (0, 0, 0))
            paste_x = (size[0] - new_width) // 2
            paste_y = (size[1] - new_height) // 2
            full_img.paste(resized_img, (paste_x, paste_y))
        else:
            full_img = resized_img.crop((left, top, right, bottom))
        
        full_img.save(temp_path)
        
        # Create base image clip
        img_clip = ImageClip(str(temp_path)).with_duration(duration)
        
        # Apply animation using resize transform (compatible with MoviePy 2.2.1)
        try:
            if animation_type == "zoom":
                # Zoom in effect using resize
                def make_frame(t):
                    # Progressive zoom from 100% to 120% over duration
                    scale = 1.0 + (0.2 * t / duration)
                    return img_clip.get_frame(t)
                
                # Apply gradual zoom by changing the clip size over time
                zoom_clip = img_clip.resized(lambda t: 1.0 + (0.2 * t / duration))
                return zoom_clip
                
            elif animation_type == "pan":
                # Pan effect using position changes
                def pos_func(t):
                    # Pan from left to center over duration
                    x_offset = int(50 * (1 - t / duration))  # Start 50px left, end at center
                    return ('center', 'center')
                
                pan_clip = img_clip.with_position(pos_func)
                return pan_clip
                
            elif animation_type == "fade":
                # Simple fade in effect
                if hasattr(img_clip, 'fadein'):
                    fade_clip = img_clip.fadein(duration * 0.3)  # Fade in over first 30%
                    return fade_clip
                else:
                    return img_clip
                    
        except Exception as e:
            # If animation fails, return basic clip
            st.warning(f"Animation failed, using static image: {str(e)}")
            
        return img_clip
    
    def generate_voiceover(self, text: str, lang: str = 'en') -> str:
        """Generate voiceover using TTS"""
        if self.progress_tracker:
            self.progress_tracker.update("Generating voiceover", f"Language: {lang}")
        
        tts = gTTS(text=text, lang=lang, slow=False)
        audio_path = self.output_dir / f"voiceover_{hash(text)}.mp3"
        tts.save(str(audio_path))
        return str(audio_path)
    
    def create_frame_clip(self, images: List[Image.Image], text: str, expand_text: bool = True, 
                         animation_type: str = "zoom") -> CompositeVideoClip:
        """Create a single frame combining images, text, and audio"""
        if self.progress_tracker:
            self.progress_tracker.update("Creating video frame", f"Images: {len(images)}")
        
        if expand_text:
            text = TextExpander.expand_text(text)
        
        # Generate voiceover
        audio_path = self.generate_voiceover(text)
        audio_clip = AudioFileClip(audio_path)
        duration = audio_clip.duration
        
        clips = []
        
        # Add images as full-screen background with animation
        if images:
            # Use first image as main background with animation
            main_img_clip = self.create_image_clip_with_animation(images[0], duration, animation_type)
            clips.append(main_img_clip)
            
            # If multiple images, create a slideshow effect within the frame
            if len(images) > 1:
                img_duration = duration / len(images)
                for i, img in enumerate(images[1:], 1):  # Skip first image (already used)
                    img_clip = self.create_image_clip_with_animation(img, img_duration, animation_type)
                    img_clip = img_clip.with_start(i * img_duration)
                    clips.append(img_clip)
        else:
            # Create animated dark background if no images
            bg_clip = ColorClip(size=(1920, 1080), color=(20, 20, 40)).with_duration(duration)
            clips.append(bg_clip)
        
        # Add semi-transparent overlay for better text visibility
        #overlay = ColorClip(size=(1920, 100), color=(0, 0, 0))
        #overlay = overlay.with_opacity(0.7).with_duration(duration).with_position(('center', 'bottom'))
        #clips.append(overlay)
        
        # Add text at bottom with better positioning
        text_clip = self.create_text_clip(text, duration)
        text_width = text_clip.w
        text_height = text_clip.h
        screen_width = 1920
        screen_height = 1080

        y_position = screen_height - text_height - 20  # 20px padding from bottom

        # Scrolling left to right
        def marquee_position(t):
            x = screen_width - (t * (screen_width + text_width) / duration)
            return (x, y_position)

        text_clip = text_clip.with_position(marquee_position)
        clips.append(text_clip)
        
        # Compose final clip
        final_clip = CompositeVideoClip(clips)
        final_clip = final_clip.with_audio(audio_clip)
        
        return final_clip
    
    def add_cross_dissolve_transitions(self, clips: List[CompositeVideoClip], transition_duration: float = 1.0) -> CompositeVideoClip:
        """Add cross-dissolve transitions between clips using overlap"""
        if len(clips) <= 1:
            return clips[0] if clips else None
        
        if self.progress_tracker:
            self.progress_tracker.update("Adding transitions", f"Transition duration: {transition_duration}s")
        
        try:
            # Calculate total duration and positions for overlapping clips
            final_clips = []
            current_time = 0
            
            for i, clip in enumerate(clips):
                if i == 0:
                    # First clip - just add fade in if available
                    if hasattr(clip, 'fadein'):
                        clip_with_fade = clip.fadein(0.5)
                        final_clips.append(clip_with_fade.with_start(current_time))
                    else:
                        final_clips.append(clip.with_start(current_time))
                    current_time += clip.duration - transition_duration
                else:
                    # Subsequent clips - add with overlap for cross-dissolve effect
                    if hasattr(clip, 'fadein') and hasattr(clips[i-1], 'fadeout'):
                        # Add fade in to current clip
                        clip_with_fade = clip.fadein(transition_duration)
                        final_clips.append(clip_with_fade.with_start(current_time))
                        
                        # Add fade out to previous clip (modify the last added clip)
                        if len(final_clips) >= 2:
                            prev_clip = final_clips[-2]
                            prev_clip_faded = prev_clip.fadeout(transition_duration)
                            final_clips[-2] = prev_clip_faded
                    else:
                        final_clips.append(clip.with_start(current_time))
                    
                    if i < len(clips) - 1:  # Not the last clip
                        current_time += clip.duration - transition_duration
                    else:  # Last clip
                        if hasattr(clip, 'fadeout'):
                            clip_with_fade = final_clips[-1].fadeout(0.5)
                            final_clips[-1] = clip_with_fade
            
            # Compose all clips together
            final_video = CompositeVideoClip(final_clips)
            return final_video
            
        except Exception as e:
            st.warning(f"Transition effects failed: {str(e)}. Using simple concatenation.")
            # Fallback to simple concatenation
            return concatenate_videoclips(clips)

def main():
    st.set_page_config(page_title="Document to Video Generator", layout="wide")
    
    st.header("Document to Video Generator")
    st.markdown("Upload documents and create engaging videos with animations and voiceovers!")
    
    # Initialize session state
    if 'extracted_content' not in st.session_state:
        st.session_state.extracted_content = None
    if 'matches' not in st.session_state:
        st.session_state.matches = []
    
    # File upload section
    st.subheader("1. Upload Document")
    uploaded_file = st.file_uploader(
        "Choose a document file",
        type=['pdf', 'doc', 'docx', 'ppt', 'pptx'],
        help="Supported formats: PDF, DOC, DOCX, PPT, PPTX"
    )
    
    if uploaded_file:
        with st.spinner("Processing document..."):
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_path = tmp_file.name
            
            try:
                # Extract content based on file type
                file_ext = uploaded_file.name.split('.')[-1].lower()
                
                if file_ext == 'pdf':
                    texts, images = DocumentProcessor.extract_from_pdf(tmp_path)
                elif file_ext in ['doc', 'docx']:
                    texts, images = DocumentProcessor.extract_from_docx(tmp_path)
                elif file_ext in ['ppt', 'pptx']:
                    texts, images = DocumentProcessor.extract_from_pptx(tmp_path)
                
                st.session_state.extracted_content = {
                    'texts': texts,
                    'images': images,
                    'filename': uploaded_file.name
                }
                
                st.success(f"✅ Extracted {len(texts)} text segments and {len(images)} images")
                
            except Exception as e:
                st.error(f"Error processing document: {str(e)}")
            finally:
                # Clean up temp file
                os.unlink(tmp_path)
    
    # Content preview and matching section
    if st.session_state.extracted_content:
        st.subheader("2. Preview Extracted Content")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("📝 Text Segments")
            texts = st.session_state.extracted_content['texts']
            for i, text in enumerate(texts):
                with st.expander(f"Text {i+1}"):
                    st.write(text[:200] + "..." if len(text) > 200 else text)
        
        with col2:
            st.markdown("🖼️ Images")
            images = st.session_state.extracted_content['images']
            for i, img in enumerate(images):
                with st.expander(f"Image {i+1}"):
                    st.image(img, width=300)
        
        # Matching interface
        st.subheader("3. Match Images with Text")
        
        with st.form("matching_form"):
            st.write("Create combinations of images and text for video frames:")
            
            # Add match
            col1, col2, col3 = st.columns([2, 2, 1])
            
            with col1:
                selected_images = st.multiselect(
                    "Select Images",
                    options=list(range(len(images))),
                    format_func=lambda x: f"Image {x+1}",
                    key="selected_images"
                )
            
            with col2:
                selected_texts = st.multiselect(
                    "Select Text Segments",
                    options=list(range(len(texts))),
                    format_func=lambda x: f"Text {x+1}",
                    key="selected_texts"
                )
            
            with col3:
                if st.form_submit_button("Add Match"):
                    if selected_images or selected_texts:
                        match = {
                            'images': selected_images,
                            'texts': selected_texts,
                            'combined_text': ' '.join([texts[i] for i in selected_texts])
                        }
                        st.session_state.matches.append(match)
                        st.success("Match added!")
        
        # Display current matches
        if st.session_state.matches:
            st.markdown("Current Matches")
            for i, match in enumerate(st.session_state.matches):
                with st.expander(f"Frame {i+1}"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write(f"**Images:** {[f'Image {j+1}' for j in match['images']]}")
                        for img_idx in match['images']:
                            if img_idx < len(images):
                                st.image(images[img_idx], width=200)
                    
                    with col2:
                        st.write(f"**Texts:** {[f'Text {j+1}' for j in match['texts']]}")
                        st.write(match['combined_text'][:300] + "..." if len(match['combined_text']) > 300 else match['combined_text'])
                    
                    if st.button(f"Remove Frame {i+1}", key=f"remove_{i}"):
                        st.session_state.matches.pop(i)
                        st.rerun()
        
        # Video generation section
        if st.session_state.matches:
            st.subheader("4. Generate Video")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                expand_text = st.checkbox("Expand text for detailed explanation", value=True)
                add_transitions = st.checkbox("Add cross-dissolve transitions", value=True)
            
            with col2:
                voice_lang = st.selectbox("Voice Language", ['en', 'es', 'fr', 'de', 'it'], index=0)
                animation_type = st.selectbox("Animation Type", ['zoom', 'pan', 'fade', 'none'], index=0)
            
            with col3:
                if add_transitions:
                    transition_duration = st.slider("Transition Duration (seconds)", 0.5, 3.0, 1.0, 0.5)
                else:
                    transition_duration = 0.0
            
            if st.button("🎬 Generate Video", type="primary"):
                # More accurate step calculation
                num_matches = len(st.session_state.matches)
                steps_per_frame = 5  # More conservative estimate
                final_steps = 3  # Initialization, assembly, rendering
                total_steps = (num_matches * steps_per_frame) + final_steps
                
                progress_tracker = ProgressTracker(total_steps)
                
                try:
                    # Create video generator
                    video_gen = VideoGenerator("temp_video")
                    video_gen.set_progress_tracker(progress_tracker)
                    
                    progress_tracker.set_progress(0.05, "Initializing video generation", "Setting up temporary files...")
                    
                    # Create clips for each match
                    clips = []
                    for i, match in enumerate(st.session_state.matches):
                        # Update progress based on frame processing
                        frame_progress = 0.1 + (0.7 * i / num_matches)  # Frames take 70% of total time
                        progress_tracker.set_progress(frame_progress, f"Processing frame {i+1}", f"Frame {i+1} of {num_matches}")
                        
                        frame_images = [images[j] for j in match['images'] if j < len(images)]
                        frame_clip = video_gen.create_frame_clip(
                            frame_images, 
                            match['combined_text'], 
                            expand_text,
                            animation_type if animation_type != 'none' else 'zoom'
                        )
                        clips.append(frame_clip)
                    
                    progress_tracker.set_progress(0.85, "Assembling video", "Combining all frames...")
                    
                    # Add transitions or concatenate clips
                    if add_transitions and len(clips) > 1:
                        final_video = video_gen.add_cross_dissolve_transitions(clips, transition_duration)
                    else:
                        final_video = concatenate_videoclips(clips)
                    
                    progress_tracker.set_progress(0.95, "Rendering video", "Encoding final video file...")
                    
                    # Save video - MoviePy 2.2.1 compatible
                    output_path = "generated_video.mp4"
                    final_video.write_videofile(
                        output_path,
                        fps=24,
                        codec='libx264',
                        audio_codec='aac',
                        logger='bar'
                    )
                    
                    progress_tracker.complete()
                    
                    st.success("🎉 Video generated successfully!")
                    
                    # Provide download link
                    with open(output_path, 'rb') as f:
                        video_bytes = f.read()
                    
                    st.download_button(
                        label="📥 Download Video",
                        data=video_bytes,
                        file_name=f"video_{st.session_state.extracted_content['filename']}.mp4",
                        mime="video/mp4"
                    )
                    
                    # Display video preview
                    st.video(output_path)
                    
                    # Clean up clips to free memory
                    for clip in clips:
                        if hasattr(clip, 'close'):
                            clip.close()
                    if hasattr(final_video, 'close'):
                        final_video.close()
                    
                    # Clean up temporary files
                    temp_files = list(Path("temp_video").glob("*"))
                    for temp_file in temp_files:
                        try:
                            temp_file.unlink()
                        except:
                            pass
                        
                except Exception as e:
                    st.error(f"Error generating video: {str(e)}")
                    st.error("Common issues:")
                    st.markdown("""
                    - **Font issues**: Try disabling text overlays
                    - **Animation issues**: Set animation to 'none'
                    - **Transition issues**: Disable transitions
                    - **Memory issues**: Reduce number of frames or image quality
                    """)

if __name__ == "__main__":
    main()