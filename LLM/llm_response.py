import requests
import json
import docx
import PyPDF2
import pandas as pd
from pptx import Presentation
from PIL import Image
import io
import base64
import re

def read_docx(file_path):
    """Extract text and images from DOCX file"""
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    
    # Extract images
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                img_data = rel.target_part.blob
                img = Image.open(io.BytesIO(img_data))
                images.append(img)
            except Exception as e:
                print(f"Error extracting image: {e}")
    
    return content, images

def read_pdf(file_path):
    """Extract text from PDF file"""
    content = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfFileReader(file)
        for page_num in range(reader.numPages):
            page = reader.getPage(page_num)
            content += page.extract_text() + "\n"
    
    # Note: PDF image extraction is complex and would require additional libraries like pymupdf
    # For now, returning empty images list
    images = []
    return content, images

def read_pptx(file_path):
    """Extract text and images from PPTX file"""
    prs = Presentation(file_path)
    content = ""
    images = []
    
    for slide in prs.slides:
        # Extract text
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
        
        # Extract images
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture shape type
                try:
                    img_data = shape.image.blob
                    img = Image.open(io.BytesIO(img_data))
                    images.append(img)
                except Exception as e:
                    print(f"Error extracting image from slide: {e}")
    
    return content, images

def process_file(file_path, file_ext):
    """Process file based on extension and return content and images"""
    try:
        if file_ext == 'pdf':
            texts, images = read_pdf(file_path)
        elif file_ext in ['doc', 'docx']:
            texts, images = read_docx(file_path)
        elif file_ext in ['ppt', 'pptx']:
            texts, images = read_pptx(file_path)
        else:
            return f"Unsupported file format: {file_ext}", []
    except Exception as e:
        return f"Error processing document: {str(e)}", []
    
    return texts, images

def get_llm_summary(file_path, file_ext):
    """Get LLM summary and return structured dataframe"""
    content, images = process_file(file_path, file_ext)
    
    if isinstance(content, str) and content.startswith("Error"):
        return pd.DataFrame({'images': [], 'description': [content]})
    
    # Combine the file content with the prompt
    combined_prompt = f"""{content}

You are a trainer and the attached document is the training guide. Create a step-by-step tutorial breakdown with each distinct step called out in details for end user understanding. 

Format your response as follows:
Step 1: [Detailed description of first step]
Step 2: [Detailed description of second step]
Step 3: [Detailed description of third step]
...and so on

Each step should be comprehensive and self-contained. Do not use any special characters or formatting instructions for narrator. Focus on clear, actionable instructions."""

    # API request setup
    url = "http://127.0.0.1:11434/api/generate"
    data = {
        "model": "qwen:latest",
        "prompt": combined_prompt
    }

    try:
        # Send the request
        response = requests.post(url, json=data)
        
        # Process the response
        raw_responses = response.text.splitlines()
        single_line_response = ""
        for raw_response in raw_responses:
            try:
                json_response = json.loads(raw_response)
                single_line_response += json_response["response"]
            except json.JSONDecodeError as e:
                print("JSONDecodeError:", e)
        
        # Parse steps from the response
        steps = parse_steps_from_response(single_line_response)
        
        # Create dataframe
        df = create_dataframe_with_images(steps, images)
        
        return df
        
    except Exception as e:
        return pd.DataFrame({'images': [], 'description': [f"Error getting LLM response: {str(e)}"]})

def parse_steps_from_response(response_text):
    """Parse individual steps from LLM response"""
    steps = []
    
    # Split by step patterns
    step_pattern = r'Step\s*\d+\s*:?\s*(.+?)(?=Step\s*\d+|$)'
    matches = re.findall(step_pattern, response_text, re.DOTALL | re.IGNORECASE)
    
    if matches:
        steps = [match.strip() for match in matches]
    else:
        # Fallback: split by sentences if no step pattern found
        sentences = response_text.split('.')
        steps = [sentence.strip() for sentence in sentences if sentence.strip()]
    
    return steps

def create_dataframe_with_images(steps, images):
    """Create dataframe with images and corresponding step descriptions"""
    df_data = {'images': [], 'description': []}
    
    # If we have images, pair them with steps
    if images and steps:
        # Distribute images across steps
        for i, step in enumerate(steps):
            if i < len(images):
                df_data['images'].append(images[i])
            else:
                # If more steps than images, reuse the last image or use None
                df_data['images'].append(images[-1] if images else None)
            df_data['description'].append(step)
    
    # If no images but have steps, create entries with None images
    elif steps:
        for step in steps:
            df_data['images'].append(None)
            df_data['description'].append(step)
    
    # If no steps, create a single entry
    else:
        df_data['images'].append(images[0] if images else None)
        df_data['description'].append("No structured steps found in the document.")
    
    return pd.DataFrame(df_data)

def image_to_base64(image):
    """Convert PIL Image to base64 string for storage/display"""
    if image is None:
        return None
    
    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()
    return img_str

def base64_to_image(base64_str):
    """Convert base64 string back to PIL Image"""
    if base64_str is None:
        return None
    
    img_data = base64.b64decode(base64_str)
    image = Image.open(io.BytesIO(img_data))
    return image