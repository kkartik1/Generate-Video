import streamlit as st
#import pandas as pd
from typing import List, Dict, Tuple, Optional
import tempfile
import os
from pathlib import Path
import base64
from io import BytesIO
import json
import time
import textwrap
import llm_response as llm
import re
import pandas as pd

# Document processing imports
try:
    import fitz  # PyMuPDF for PDF
    import docx  # python-docx for Word documents
    from pptx import Presentation  # python-pptx for PowerPoint
    import mammoth  # For better DOC/DOCX text extraction
    from PIL import Image
    import pytesseract  # For OCR if needed
except ImportError as e:
    st.error(f"Missing required library: {e}")
    st.stop()

# Video generation imports
try:
    # MoviePy 2.2.1 compatible imports
    from moviepy import VideoFileClip, AudioFileClip, ImageClip, TextClip, ColorClip, CompositeVideoClip, concatenate_videoclips
    import cv2
    import numpy as np
    from gtts import gTTS  # Google Text-to-Speech
    
    # Test MoviePy installation
    test_clip = ColorClip(size=(100, 100), color=(0, 0, 0), duration=0.1)
    test_clip.close()
    
except ImportError as e:
    st.error(f"Missing video processing library: {e}")
    st.markdown("""
    ### 🔧 Installation Instructions
    
    For MoviePy 2.2.1 compatibility:
    
    **Recommended installation:**
    ```bash
    pip install "moviepy==2.2.1"
    pip install imageio-ffmpeg
    ```
    
    **Alternative methods:**
    ```bash
    # Method 1 - With conda
    conda install -c conda-forge moviepy
    
    # Method 2 - Latest version
    pip install moviepy
    ```
    
    After installation, restart your Python environment/kernel.
    """)
    st.stop()
except Exception as e:
    st.error(f"MoviePy installation issue: {e}")
    st.markdown("""
    ### Additional Requirements
    
    MoviePy also requires:
    - **FFmpeg** (for video processing)
    - **ImageIO** (for image handling)
    
    **To install FFmpeg:**
    
    **Windows:**
    ```bash
    # Using chocolatey
    choco install ffmpeg
    
    # Or download from https://ffmpeg.org/download.html
    ```
    
    **macOS:**
    ```bash
    brew install ffmpeg
    ```
    
    **Linux (Ubuntu/Debian):**
    ```bash
    sudo apt update
    sudo apt install ffmpeg
    ```
    
    **Alternative: Install imageio-ffmpeg**
    ```bash
    pip install imageio-ffmpeg
    ```
    """)
    st.stop()

class DocumentProcessor:
    """Handle document text and image extraction"""
    
    @staticmethod
    def extract_from_pdf(file_path: str) -> Tuple[List[str], List[Image.Image]]:
        """Extract text and images from PDF"""
        doc = fitz.open(file_path)
        texts = []
        images = []
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            
            # Extract text
            text = page.get_text()
            if text.strip():
                texts.append(text.strip())
            
            # Extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    img_data = pix.tobytes("png")
                    img_pil = Image.open(BytesIO(img_data))
                    images.append(img_pil)
                pix = None
        
        doc.close()
        return texts, images
    
    @staticmethod
    def extract_from_docx(file_path: str) -> Tuple[List[str], List[Image.Image]]:
        """Extract text and images from DOCX"""
        doc = docx.Document(file_path)
        texts = []
        images = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                texts.append(paragraph.text.strip())
        
        # Extract images from document parts
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_data = rel.target_part.blob
                    img_pil = Image.open(BytesIO(img_data))
                    images.append(img_pil)
                except Exception:
                    continue
        
        return texts, images
    
    @staticmethod
    def extract_from_pptx(file_path: str) -> Tuple[List[str], List[Image.Image]]:
        """Extract text and images from PPTX"""
        prs = Presentation(file_path)
        texts = []
        images = []
        
        for slide in prs.slides:
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract images
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        img_data = image.blob
                        img_pil = Image.open(BytesIO(img_data))
                        images.append(img_pil)
                    except Exception:
                        continue
            
            if slide_text:
                texts.append(" ".join(slide_text))
        
        return texts, images


class ProgressTracker:
    """Track and display progress for video generation"""
    
    def __init__(self, total_steps: int):
        self.total_steps = max(1, total_steps)  # Ensure total_steps is at least 1
        self.current_step = 0
        self.progress_bar = st.progress(0)
        self.status_text = st.empty()
        self.detail_text = st.empty()
        
    def update(self, step_name: str, detail: str = ""):
        """Update progress bar and status"""
        self.current_step += 1
        # Ensure progress never exceeds 1.0
        progress = min(self.current_step / self.total_steps, 1.0)
        self.progress_bar.progress(progress)
        self.status_text.text(f"Step {min(self.current_step, self.total_steps)}/{self.total_steps}: {step_name}")
        if detail:
            self.detail_text.text(detail)
        
    def set_progress(self, progress_value: float, step_name: str, detail: str = ""):
        """Set progress directly with a value between 0.0 and 1.0"""
        # Ensure progress is within valid bounds
        progress = max(0.0, min(progress_value, 1.0))
        self.progress_bar.progress(progress)
        self.status_text.text(step_name)
        if detail:
            self.detail_text.text(detail)
        
    def complete(self):
        """Mark as complete"""
        self.progress_bar.progress(1.0)
        self.status_text.text("✅ Video generation complete!")
        self.detail_text.text("")

class VideoGenerator:
    """Generate video from matched content"""
    
    def __init__(self, output_dir: str):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.progress_tracker = None
    
    def set_progress_tracker(self, tracker: ProgressTracker):
        """Set progress tracker for this generator"""
        self.progress_tracker = tracker
    
    def create_text_clip(self, text: str, duration: float, size: Tuple[int, int] = (1600, 800)) -> TextClip:
        """Create text clip with improved styling"""
        if self.progress_tracker:
            self.progress_tracker.update("Creating text overlay", f"Processing: {text[:50]}...")
        
        # Create base text clip with better visibility - MoviePy 2.2.1 compatible
        try:
            # Try common system fonts in order of preference
            fonts_to_try = ['Arial', 'DejaVu-Sans', 'Liberation-Sans']
            font_used = None
            wrapped_text = textwrap.fill(text, width=70)  # Wrap text to approx. 70 chars per line
            for font in fonts_to_try:
                try:
                    txt_clip = TextClip(
                        text=wrapped_text,
                        font_size=40,
                        color='white', 
                        font=font,
                        stroke_color='black',
                        stroke_width=2,
                        method='caption',
                        size=(1600, None)  # Width defined; height auto-adjusted
                    ).with_duration(duration)
                    font_used = font
                    break
                except:
                    continue
            
            # If no specific font works, use default
            if font_used is None:
                txt_clip = TextClip(
                    text=wrapped_text,
                    font_size=40,
                    color='white',
                    stroke_color='black',
                    stroke_width=2,
                    method='caption',
                    size=(1600, None)
                ).with_duration(duration)
                
        except Exception as e:
            # Fallback to simplest TextClip creation
            txt_clip = TextClip(
                text=wrapped_text,
                font_size=30,
                color='white'
            ).with_duration(duration)
        
        return txt_clip
    
    def create_image_clip_with_animation(self, image: Image.Image, duration: float, 
                                       size: Tuple[int, int] = (1600, 800)) -> ImageClip:
        """Create image clip with animation effects using resize transforms"""
       
        # Save image temporarily
        temp_path = self.output_dir / f"temp_img_{id(image)}.png"
        
        # Resize image to full screen while maintaining aspect ratio
        img_ratio = image.width / image.height
        screen_ratio = size[0] / size[1]
        
        if img_ratio > screen_ratio:
            # Image is wider - fit to height
            new_height = size[1]
            new_width = int(new_height * img_ratio)
        else:
            # Image is taller - fit to width
            new_width = size[0]
            new_height = int(new_width / img_ratio)
        
        # Resize and center crop to full screen
        resized_img = image.resize((new_width, new_height), Image.LANCZOS)
        
        # Create full screen image by center cropping
        left = (new_width - size[0]) // 2
        top = (new_height - size[1]) // 2
        right = left + size[0]
        bottom = top + size[1]
        
        if left < 0 or top < 0:
            # If image is smaller than screen, create background and paste
            full_img = Image.new('RGB', size, (0, 0, 0))
            paste_x = (size[0] - new_width) // 2
            paste_y = (size[1] - new_height) // 2
            full_img.paste(resized_img, (paste_x, paste_y))
        else:
            full_img = resized_img.crop((left, top, right, bottom))
        
        full_img.save(temp_path)
        
        # Create base image clip
        img_clip = ImageClip(str(temp_path)).with_duration(duration)
        
        return img_clip
    
    def generate_voiceover(self, text: str, lang: str = 'en') -> str:
        """Generate voiceover using TTS"""
        if self.progress_tracker:
            self.progress_tracker.update("Generating voiceover", f"Language: {lang}")
        
        tts = gTTS(text=text, lang=lang, slow=False)
        audio_path = self.output_dir / f"voiceover_{hash(text)}.mp3"
        tts.save(str(audio_path))
        return str(audio_path)
    
    def create_frame_clip(self, images: List[Image.Image], text: str) -> CompositeVideoClip:
        """Create a single frame combining images, text, and audio"""
        if self.progress_tracker:
            self.progress_tracker.update("Creating video frame", f"Images: {len(images)}")
        
        # Generate voiceover
        audio_path = self.generate_voiceover(text)
        audio_clip = AudioFileClip(audio_path)
        duration = audio_clip.duration
        
        clips = []
        
        # Add images as full-screen background with animation
        if images:
            # Use first image as main background with animation
            main_img_clip = self.create_image_clip_with_animation(images[0], duration)
            clips.append(main_img_clip)
            
            # If multiple images, create a slideshow effect within the frame
            if len(images) > 1:
                img_duration = duration / len(images)
                for i, img in enumerate(images[1:], 1):  # Skip first image (already used)
                    img_clip = self.create_image_clip_with_animation(img, img_duration)
                    img_clip = img_clip.with_start(i * img_duration)
                    clips.append(img_clip)
        else:
            # Create animated dark background if no images
            bg_clip = ColorClip(size=(1600, 800), color=(20, 20, 40)).with_duration(duration)
            clips.append(bg_clip)
        
        # Add text at bottom with better positioning
        text_clip = self.create_text_clip(text, duration)
        text_width = text_clip.w
        text_height = text_clip.h
        screen_width = 1600
        screen_height = 800

        y_position = screen_height - text_height - 20  # 20px padding from bottom

        text_clip = text_clip.with_position(('center', y_position))
        clips.append(text_clip)
        
        # Compose final clip
        final_clip = CompositeVideoClip(clips)
        final_clip = final_clip.with_audio(audio_clip)
        
        return final_clip
    
    def add_cross_dissolve_transitions(self, clips: List[CompositeVideoClip], transition_duration: float = 1.0) -> CompositeVideoClip:
        """Add cross-dissolve transitions between clips with proper audio handling - FIXED VERSION"""
        if len(clips) <= 1:
            return clips[0] if clips else None
        
        if self.progress_tracker:
            self.progress_tracker.update("Adding transitions", f"Transition duration: {transition_duration}s")
        
        try:
            # Reduce transition duration if it's too long compared to clip durations
            min_clip_duration = min(clip.duration for clip in clips)
            max_transition = min_clip_duration / 3  # More conservative to avoid audio overlap
            effective_transition = min(transition_duration, max_transition)
            
            # Create list to hold all clips with proper timing
            final_clips = []
            current_start_time = 0
            
            for i, clip in enumerate(clips):
                if i == 0:
                    # First clip starts at time 0
                    clip_with_timing = clip.with_start(current_start_time)
                    
                    # Add fade in effect to video only, keep audio as is
                    if hasattr(clip_with_timing, 'fadein'):
                        # Create separate video and audio components
                        video_part = clip_with_timing.without_audio().fadein(0.5)
                        audio_part = clip_with_timing.audio
                        # Recombine with original audio timing
                        clip_with_timing = video_part.with_audio(audio_part)
                    
                    final_clips.append(clip_with_timing)
                    
                    # Next clip starts after this one, minus transition duration
                    # But audio should not overlap, so we calculate differently
                    current_start_time += clip.duration - effective_transition
                    
                else:
                    # For subsequent clips, handle video and audio separately
                    original_audio = clip.audio
                    original_video = clip.without_audio()
                    
                    # Video part starts with overlap for visual transition
                    video_part = original_video.with_start(current_start_time)
                    if hasattr(video_part, 'fadein'):
                        video_part = video_part.fadein(effective_transition)
                    
                    # Audio part starts AFTER the previous audio ends (no overlap)
                    # Calculate when previous audio actually ends
                    prev_clip_audio_end = final_clips[-1].audio.end if hasattr(final_clips[-1], 'audio') and final_clips[-1].audio else current_start_time
                    audio_start_time = max(current_start_time, prev_clip_audio_end)
                    
                    # Ensure audio doesn't start before its video
                    audio_start_time = max(audio_start_time, current_start_time)
                    
                    audio_part = original_audio.with_start(audio_start_time)
                    
                    # Add fade out to previous clip's video only
                    if len(final_clips) > 0:
                        prev_clip = final_clips[-1]
                        if hasattr(prev_clip, 'fadeout'):
                            # Separate video and audio of previous clip
                            prev_video = prev_clip.without_audio()
                            prev_audio = prev_clip.audio
                            
                            # Apply fadeout only to video
                            prev_video_faded = prev_video.fadeout(effective_transition)
                            
                            # Recombine with original audio (no fade on audio)
                            final_clips[-1] = prev_video_faded.with_audio(prev_audio)
                    
                    # Combine video and audio for current clip
                    current_clip = CompositeVideoClip([video_part]).with_audio(audio_part)
                    final_clips.append(current_clip)
                    
                    # Calculate start time for next clip
                    if i < len(clips) - 1:
                        current_start_time += clip.duration - effective_transition
            
            # Add fade out to last clip's video only
            if len(final_clips) > 0:
                last_clip = final_clips[-1]
                if hasattr(last_clip, 'fadeout'):
                    last_video = last_clip.without_audio()
                    last_audio = last_clip.audio
                    last_video_faded = last_video.fadeout(0.5)
                    final_clips[-1] = last_video_faded.with_audio(last_audio)
            
            # Compose all clips together
            final_video = CompositeVideoClip(final_clips)
            return final_video
            
        except Exception as e:
            st.warning(f"Advanced transition effects failed: {str(e)}. Using simple concatenation.")
            # Fallback to simple concatenation which naturally avoids audio overlap
            return concatenate_videoclips(clips)
                
        except Exception as e:
            st.warning(f"Transition effects failed: {str(e)}. Using simple concatenation.")
            # Fallback to simple concatenation
            return concatenate_videoclips(clips)

# Function to split rows without breaking words - FIXED VERSION
def split_rows(matches_list, length=150):
    """Split text content into smaller chunks for better video pacing"""
    new_matches = []
    
    for match in matches_list:
        combined_text = match.get('combined_text', '')
        
        # If text is short enough, keep as is
        if len(combined_text) <= length:
            new_matches.append(match)
            continue
        
        # Split long text into chunks
        text_chunks = []
        remaining_text = combined_text
        
        while len(remaining_text) > length:
            # Find the last space within the length limit
            split_index = remaining_text[:length].rfind(' ')
            if split_index == -1:  # No space found, split at length
                split_index = length
            
            text_chunks.append(remaining_text[:split_index].strip())
            remaining_text = remaining_text[split_index:].strip()
        
        # Add remaining text
        if remaining_text:
            text_chunks.append(remaining_text)
        
        # Create new matches for each text chunk
        for chunk in text_chunks:
            new_match = match.copy()
            new_match['combined_text'] = chunk
            new_matches.append(new_match)
    
    return new_matches
        
def main():
    st.set_page_config(page_title="Document to Video Generator", layout="wide")
    
    st.header("Document to Video Generator")
    st.markdown("Upload documents and create engaging videos with animations and voiceovers!")
    
    # Initialize session state
    if 'extracted_content' not in st.session_state:
        st.session_state.extracted_content = None
    if 'matches' not in st.session_state:
        st.session_state.matches = []
    if 'read_content' not in st.session_state:
        st.session_state.read_content = False
        
    
    # File upload section
    st.subheader("1. Upload Document")
    uploaded_file = st.file_uploader(
        "Choose a document file",
        type=['pdf', 'doc', 'docx', 'ppt', 'pptx'],
        help="Supported formats: PDF, DOC, DOCX, PPT, PPTX"
    )
    
    if uploaded_file and not st.session_state.read_content:
        with st.spinner("Processing document..."):
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_path = tmp_file.name
            
            try:
                # Extract content based on file type
                file_ext = uploaded_file.name.split('.')[-1].lower()
                
                # call LLM to summarize data
                text_summ = llm.get_llm_summary(tmp_path, file_ext)
                print(text_summ)
                if file_ext == 'pdf':
                    texts, images = DocumentProcessor.extract_from_pdf(tmp_path)
                elif file_ext in ['doc', 'docx']:
                    texts, images = DocumentProcessor.extract_from_docx(tmp_path)
                elif file_ext in ['ppt', 'pptx']:
                    texts, images = DocumentProcessor.extract_from_pptx(tmp_path)
                
                # Split the tutorial text into steps
                texts = text_summ.split("Step")
                if len(texts) == 1:
                    texts = texts.splitlines()
                
                # Remove empty strings and strip leading/trailing whitespace
                texts = [re.sub(r'[\n\r]', ' ', step) for step in texts]
                texts = [re.sub(r'[^a-zA-Z0-9 .,!?;:\'\-"]', '', step) for step in texts]
                texts = [re.sub(r'\s+', ' ', step) for step in texts]
                texts = [step.replace("Narrator", "") for step in texts]
                texts = [step.strip() for step in texts if step.strip()]

                st.session_state.extracted_content = {
                    'texts': texts,
                    'images': images,
                    'filename': uploaded_file.name
                }
                
                st.success(f"✅ Extracted {len(texts)} text segments and {len(images)} images")
                
            except Exception as e:
                st.error(f"Error processing document: {str(e)}")
            finally:
                # Clean up temp file
                os.unlink(tmp_path)
                st.session_state.read_content = True
                
    
    # Content preview and matching section
    if st.session_state.extracted_content:
        st.subheader("2. Preview Extracted Content")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("📝 Text Segments")
            texts = st.session_state.extracted_content['texts']
            for i, text in enumerate(texts):
                with st.expander(f"Text {i+1}"):
                    st.write(text[:200] + "..." if len(text) > 200 else text)
        
        with col2:
            st.markdown("🖼️ Images")
            images = st.session_state.extracted_content['images']
            for i, img in enumerate(images):
                with st.expander(f"Image {i+1}"):
                    st.image(img, width=300)
        
        # Matching interface
        st.subheader("3. Match Images with Text")
        
        with st.form("matching_form"):
            st.write("Create combinations of images and text for video frames:")
            
            # Add match
            col1, col2, col3 = st.columns([2, 2, 1])
            
            with col1:
                selected_images = st.multiselect(
                    "Select Images",
                    options=list(range(len(images))),
                    format_func=lambda x: f"Image {x+1}",
                    key="selected_images"
                )
            
            with col2:
                selected_texts = st.multiselect(
                    "Select Text Segments",
                    options=list(range(len(texts))),
                    format_func=lambda x: f"Text {x+1}",
                    key="selected_texts"
                )
            
            with col3:
                if st.form_submit_button("Add Match"):
                    if selected_images or selected_texts:
                        match = {
                            'images': selected_images,
                            'texts': selected_texts,
                            'combined_text': ' '.join([texts[i] for i in selected_texts])
                        }
                        st.session_state.matches.append(match)
                        st.success("Match added!")
        
        # Display current matches
        if st.session_state.matches:
            st.markdown("Current Matches")
            for i, match in enumerate(st.session_state.matches):
                with st.expander(f"Frame {i+1}"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write(f"**Images:** {[f'Image {j+1}' for j in match['images']]}")
                        for img_idx in match['images']:
                            if img_idx < len(images):
                                st.image(images[img_idx], width=200)
                    
                    with col2:
                        st.write(f"**Texts:** {[f'Text {j+1}' for j in match['texts']]}")
                        st.write(match['combined_text'][:300] + "..." if len(match['combined_text']) > 300 else match['combined_text'])
                    
                    if st.button(f"Remove Frame {i+1}", key=f"remove_{i}"):
                        st.session_state.matches.pop(i)
                        st.rerun()
        
        # Video generation section
        if st.session_state.matches:
            st.subheader("4. Generate Video")
            
            # Video settings
            col1, col2 = st.columns(2)
            with col1:
                add_transitions = st.checkbox("Add cross-dissolve transitions", value=True)
            with col2:
                transition_duration = st.slider("Transition duration (seconds)", 0.5, 3.0, 1.0, 0.1)
            
            if st.button("🎬 Generate Video", type="primary"):
                # Split long text segments for better pacing
                processed_matches = split_rows(st.session_state.matches, length=150)
                
                # More accurate step calculation
                num_matches = len(processed_matches)
                print(f"Processing {num_matches} video segments")
                steps_per_frame = 3  # Conservative estimate
                final_steps = 3  # Initialization, assembly, rendering
                total_steps = (num_matches * steps_per_frame) + final_steps
                
                progress_tracker = ProgressTracker(total_steps)
                
                try:
                    # Create video generator
                    video_gen = VideoGenerator("temp_video")
                    video_gen.set_progress_tracker(progress_tracker)
                    
                    progress_tracker.set_progress(0.05, "Initializing video generation", "Setting up temporary files...")
                    
                    # Create clips for each match
                    clips = []
                    for i, match in enumerate(processed_matches):
                        # Update progress based on frame processing
                        frame_progress = 0.1 + (0.7 * i / num_matches)  # Frames take 70% of total time
                        progress_tracker.set_progress(frame_progress, f"Processing frame {i+1}", f"Frame {i+1} of {num_matches}")
                        
                        frame_images = [images[j] for j in match['images'] if j < len(images)]
                        frame_clip = video_gen.create_frame_clip(
                            frame_images, 
                            match['combined_text']
                        )
                        clips.append(frame_clip)
                    
                    progress_tracker.set_progress(0.85, "Assembling video", "Combining all frames...")
                    
                    # Add transitions or concatenate clips
                    if add_transitions and len(clips) > 1:
                        final_video = video_gen.add_cross_dissolve_transitions(clips, transition_duration)
                    else:
                        final_video = concatenate_videoclips(clips)
                    
                    progress_tracker.set_progress(0.95, "Rendering video", "Encoding final video file...")
                    
                    # Save video - MoviePy 2.2.1 compatible
                    output_path = "generated_video.mp4"
                    final_video.write_videofile(
                        output_path,
                        fps=24,
                        codec='libx264',
                        audio_codec='aac',
                        logger='bar'
                    )
                    
                    progress_tracker.complete()
                    
                    st.success("🎉 Video generated successfully!")
                    
                    # Provide download link
                    with open(output_path, 'rb') as f:
                        video_bytes = f.read()
                    
                    st.download_button(
                        label="📥 Download Video",
                        data=video_bytes,
                        file_name=f"video_{st.session_state.extracted_content['filename']}.mp4",
                        mime="video/mp4"
                    )
                    
                    # Display video preview
                    st.video(output_path)
                    
                    # Clean up clips to free memory
                    for clip in clips:
                        if hasattr(clip, 'close'):
                            clip.close()
                    if hasattr(final_video, 'close'):
                        final_video.close()
                    
                    # Clean up temporary files
                    temp_files = list(Path("temp_video").glob("*"))
                    for temp_file in temp_files:
                        try:
                            temp_file.unlink()
                        except:
                            pass
                        
                except Exception as e:
                    st.error(f"Error generating video: {str(e)}")
                    st.error("Common issues:")
                    st.markdown("""
                    - **Font issues**: Try disabling text overlays
                    - **Animation issues**: Set animation to 'none'
                    - **Transition issues**: Disable transitions
                    - **Memory issues**: Reduce number of frames or image quality
                    """)

if __name__ == "__main__":
    main()