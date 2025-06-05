import requests
import json
import docx
import PyPDF2
from pptx import Presentation

def read_docx(file_path):
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content

def read_pdf(file_path):
    content = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfFileReader(file)
        for page_num in range(reader.numPages):
            page = reader.getPage(page_num)
            content += page.extract_text() + "\n"
    return content

def read_pptx(file_path):
    prs = Presentation(file_path)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content

def process_file(file_path, file_ext):
    try:
        if file_ext == 'pdf':
            texts = read_pdf(file_path)
        elif file_ext in ['doc', 'docx']:
            texts = read_docx(file_path)
        elif file_ext in ['ppt', 'pptx']:
            texts = read_pptx(file_path)
    except Exception as e:
        return f"Error processing document: {str(e)}"
    return texts

def get_llm_summary(file_path, file_ext):
    content = process_file(file_path, file_ext)
    # Combine the file content with the prompt
    combined_prompt = f"{content}\n\nYou are a trainer and the attached document is the training guide. Use of the document to create a voiceover tutorial with each distinct steps called out in details for end user understanding. For example 'Step1 Navigate to the patient account Summary screen in eFR'. Do not use any special character or instructions for narrator."
    # API request setup
    url = "http://rn000116071:11435/api/generate"
    data = {
        "model": "qwen2.5-coder:7b",
        "prompt": combined_prompt
    }

    # Send the request
    response = requests.post(url, json=data)

    # Process the response
    raw_responses = response.text.splitlines()
    single_line_response = ""
    for raw_response in raw_responses:
        try:
            json_response = json.loads(raw_response)
            single_line_response += json_response["response"]
        except json.JSONDecodeError as e:
            print("JSONDecodeError:", e)
    return single_line_response

